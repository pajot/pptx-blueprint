import pathlib
import pptx
import re
from typing import Union, Iterable, Tuple
from pptx.shapes.base import BaseShape

_Pathlike = Union[str, pathlib.Path]


class Template:
    """Helper class for modifying pptx templates.
    """

    def __init__(self, filename: _Pathlike) -> None:
        """Initializes a Template-Modifier.

        Args:
            filename (path-like): file name or path
        """
        self._template_path = filename
        self._presentation = pptx.Presentation(filename)
        self._copy_tags_to_name()
        pass

    def replace_text(self, label: str, new_text: str) -> None:
        """Replaces text placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            text (str): new content
            scope: None, slide number, Slide object or iterable of Slide objects
        """
        slide_number, tag_name = self._parse_label(label)
        shapes = self._find_shapes(slide_number, tag_name)
        for shape in shapes:
            shape.text = new_text

    def replace_picture(self, label: str, filename: _Pathlike) -> None:
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            filename (path-like): path to an image file
        """
        pass

    def replace_table(self, label: str, data) -> None:
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            data (pandas.DataFrame): table to be inserted into the presentation
        """
        pass

    def _parse_label(self, label: str) -> Tuple[Union[int, str], str]:
        slide_number, tag_name = label.split(':')
        return int(slide_number) if slide_number != '*' else slide_number, tag_name

    def _find_shapes(self,
                     slide_number: Union[int, str],
                     tag_name: str) -> Iterable[BaseShape]:
        """Finds all shapes that match the label

        Args:
            label (str): label of the placeholder (without curly braces)
        """
        matched_shapes = []

        def _find_shapes_in_slide(slide):
            return filter(lambda shape: shape.name == f'{{{tag_name}}}', slide.shapes)

        if slide_number == '*':
            slides = self._presentation.slides
        else:
            # in label we are using 1 based indexing
            slide_index = slide_number - 1
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Can't find slide number {slide_number}.")
            slides = [self._presentation.slides[slide_index]]

        for slide in slides:
            matched_shapes.extend(_find_shapes_in_slide(slide))

        return matched_shapes

    def _get_all_shapes(self) -> Iterable[BaseShape]:
        # Do we need all the shapes? Perhaps we should filter on tags here.
        all_shapes = [shape for slide in self._presentation.slides for shape in slide.shapes]
        return all_shapes

    def _copy_tags_to_name(self) -> None:
        all_shapes = self._get_all_shapes()
        # This regex matches on tags
        regex_tag = re.compile(r'^\s*(\{\w+\})\s*$')
        for shape in all_shapes:
            # We only copy contents we recognize as tags
            if regex_tag.match(shape.text):
                shape.name = regex_tag.group(1)

    def save(self, filename: _Pathlike) -> None:
        """Saves the updated pptx to the specified filepath.

        Args:
            filename (path-like): file name or path
        """
        # TODO: make sure that the user does not override the self._template_path
        pass
